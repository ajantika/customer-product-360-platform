"""Generate a PowerPoint deck for a customer or product — uploadable to Google Slides."""
from __future__ import annotations

import io
from datetime import date

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu

# Brand colours
PURPLE     = RGBColor(0x63, 0x66, 0xF1)  # indigo-500
PURPLE_DARK= RGBColor(0x0D, 0x0B, 0x1E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY       = RGBColor(0xA7, 0x8B, 0xFA)
RED        = RGBColor(0xF8, 0x71, 0x71)
GREEN      = RGBColor(0x34, 0xD3, 0x99)
YELLOW     = RGBColor(0xFB, 0xBF, 0x24)
DARK_BG    = RGBColor(0x1E, 0x12, 0x45)

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height


def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # completely blank


def _bg(slide, color: RGBColor = PURPLE_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _txt(slide, text: str, left, top, width, height, size=18, bold=False, color=WHITE, align="left"):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color


def _rect(slide, left, top, width, height, fill_color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()


def _kpi_box(slide, label: str, value: str, left, top, width=Inches(2.4), height=Inches(1.1)):
    _rect(slide, left, top, width, height, RGBColor(0x2D, 0x1F, 0x6E))
    _txt(slide, label, left + Inches(0.12), top + Inches(0.08), width, Inches(0.3),
         size=9, color=GREY)
    _txt(slide, value, left + Inches(0.12), top + Inches(0.32), width, Inches(0.7),
         size=20, bold=True, color=WHITE)


# ─── Cover slide ──────────────────────────────────────────────────────────────

def _cover(prs: Presentation, title: str, subtitle: str, generated: str):
    slide = _blank(prs)
    _bg(slide, PURPLE_DARK)
    _rect(slide, Inches(0), Inches(0), Inches(0.35), H, PURPLE)
    _txt(slide, "CUSTOMER & PRODUCT 360", Inches(0.6), Inches(1.8), Inches(12), Inches(0.5),
         size=11, color=GREY, align="left")
    _txt(slide, title,    Inches(0.6), Inches(2.3), Inches(12), Inches(1.4), size=38, bold=True, align="left")
    _txt(slide, subtitle, Inches(0.6), Inches(3.7), Inches(10), Inches(0.6), size=16, color=GREY, align="left")
    _txt(slide, f"Generated {generated}", Inches(0.6), H - Inches(0.7), Inches(6), Inches(0.4),
         size=9, color=RGBColor(0x6B, 0x72, 0x80), align="left")


# ─── Customer 360 deck ────────────────────────────────────────────────────────

def customer_deck(
    cust: pd.Series,
    prod_df: pd.DataFrame,
    mom_df: pd.DataFrame,
    region_tbl: pd.DataFrame,
) -> bytes:
    prs = _prs()
    today = date.today().strftime("%B %d, %Y")

    status_color = {"active": GREEN, "at_risk": YELLOW, "churned": RGBColor(0x94,0xA3,0xB8)}.get(
        str(cust.get("status", "")), GREY)

    # Slide 1: Cover
    _cover(prs, str(cust["name"]), f"{cust['region']}  ·  {cust['plan_tier']}  ·  {cust['industry']}", today)

    # Slide 2: Account snapshot
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, Inches(0), Inches(0), W, Inches(1.1), RGBColor(0x2D, 0x1F, 0x6E))
    _txt(slide, f"Account Snapshot — {cust['name']}", Inches(0.4), Inches(0.22),
         Inches(10), Inches(0.65), size=22, bold=True)

    kpis = [
        ("Region",        cust.get("region", "—")),
        ("Plan Tier",     cust.get("plan_tier", "—")),
        ("MRR",           f"${float(cust.get('mrr_usd', 0)):,.0f}"),
        ("Status",        str(cust.get("status", "—")).replace("_", " ").title()),
        ("Contract End",  pd.Timestamp(cust["contract_end_date"]).strftime("%b %Y") if "contract_end_date" in cust.index else "—"),
    ]
    for i, (lbl, val) in enumerate(kpis):
        _kpi_box(slide, lbl, val, Inches(0.4 + i * 2.55), Inches(1.4))

    _txt(slide, "Acquisition Channel", Inches(0.4), Inches(2.85), Inches(4), Inches(0.35), size=10, color=GREY)
    _txt(slide, f"📣  {cust.get('marketing_campaign', '—')}", Inches(0.4), Inches(3.15),
         Inches(6), Inches(0.45), size=14, bold=True)

    # Slide 3: Product utilization
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, Inches(0), Inches(0), W, Inches(1.1), RGBColor(0x2D, 0x1F, 0x6E))
    _txt(slide, "Product Utilization", Inches(0.4), Inches(0.22), Inches(10), Inches(0.65), size=22, bold=True)

    if not prod_df.empty:
        headers = ["Product", "Category", "Usage", "Plan Limit", "Utilization"]
        col_w   = [Inches(2.5), Inches(1.5), Inches(2.2), Inches(2.2), Inches(2.0)]
        row_h   = Inches(0.5)
        top0    = Inches(1.3)
        lefts   = [Inches(0.3)]
        for w in col_w[:-1]:
            lefts.append(lefts[-1] + w)

        for ci, (h, lft, ww) in enumerate(zip(headers, lefts, col_w)):
            _rect(slide, lft, top0, ww, row_h, RGBColor(0x3B, 0x2A, 0x8A))
            _txt(slide, h, lft + Inches(0.08), top0 + Inches(0.1), ww, row_h, size=10, bold=True, color=GREY)

        for ri, (_, row) in enumerate(prod_df.iterrows()):
            row_top = top0 + row_h * (ri + 1)
            bg = RGBColor(0x1A, 0x10, 0x40) if ri % 2 == 0 else RGBColor(0x22, 0x15, 0x55)
            for lft, ww in zip(lefts, col_w):
                _rect(slide, lft, row_top, ww, row_h, bg)
            util = float(row["utilization_pct"])
            util_color = RED if util > 100 else (YELLOW if util < 50 else GREEN)
            vals = [
                row["name"], row["category"],
                f"{row['usage']:,.1f} {row['unit']}",
                f"{row['plan_limit']:,.1f} {row['unit']}",
                f"{util:.0f}%",
            ]
            for vi, (val, lft, ww) in enumerate(zip(vals, lefts, col_w)):
                c = util_color if vi == 4 else WHITE
                _txt(slide, str(val), lft + Inches(0.08), row_top + Inches(0.12), ww, row_h, size=11, color=c)

    # Slide 4: Pricing region table
    if not region_tbl.empty:
        slide = _blank(prs)
        _bg(slide)
        _rect(slide, Inches(0), Inches(0), W, Inches(1.1), RGBColor(0x2D, 0x1F, 0x6E))
        _txt(slide, "Usage by Pricing Region", Inches(0.4), Inches(0.22), Inches(10), Inches(0.65), size=22, bold=True)
        _txt(slide, "High price-multiplier regions with high usage share = cost concentration risk",
             Inches(0.4), Inches(1.15), Inches(12), Inches(0.4), size=10, color=GREY)

        headers = ["Region", "Usage %", "Price Multiplier", "Weighted MRR", "Alert"]
        col_w   = [Inches(3.0), Inches(1.5), Inches(2.2), Inches(2.2), Inches(2.0)]
        row_h   = Inches(0.52)
        top0    = Inches(1.65)
        lefts   = [Inches(0.3)]
        for w in col_w[:-1]:
            lefts.append(lefts[-1] + w)

        for ci, (h, lft, ww) in enumerate(zip(headers, lefts, col_w)):
            _rect(slide, lft, top0, ww, row_h, RGBColor(0x3B, 0x2A, 0x8A))
            _txt(slide, h, lft + Inches(0.08), top0 + Inches(0.12), ww, row_h, size=10, bold=True, color=GREY)

        for ri, (_, row) in enumerate(region_tbl.iterrows()):
            row_top = top0 + row_h * (ri + 1)
            bg = RGBColor(0x1A, 0x10, 0x40) if ri % 2 == 0 else RGBColor(0x22, 0x15, 0x55)
            for lft, ww in zip(lefts, col_w):
                _rect(slide, lft, row_top, ww, row_h, bg)
            vals = [
                str(row.get("Region", row.get("name", ""))),
                str(row.get("Usage %", "")),
                str(row.get("Price mult.", "")),
                str(row.get("Weighted MRR", "")),
                str(row.get("Alert", "")),
            ]
            for val, lft, ww in zip(vals, lefts, col_w):
                _txt(slide, val, lft + Inches(0.08), row_top + Inches(0.13), ww, row_h, size=11)

    # Slide 5: MoM trend (as a table since embedding chart images needs kaleido)
    if not mom_df.empty:
        slide = _blank(prs)
        _bg(slide)
        _rect(slide, Inches(0), Inches(0), W, Inches(1.1), RGBColor(0x2D, 0x1F, 0x6E))
        _txt(slide, "Month-over-Month Usage Trend", Inches(0.4), Inches(0.22), Inches(10), Inches(0.65), size=22, bold=True)
        recent = mom_df.tail(12).copy()
        recent["month_label"] = recent["month"].dt.strftime("%b %Y")
        col_w = [Inches(2.2), Inches(2.2)]
        row_h = Inches(0.42)
        top0  = Inches(1.3)
        for ri, (_, row) in enumerate(recent.iterrows()):
            col = ri % 2
            rr  = ri // 2
            lft = Inches(0.3) + col * Inches(6.5)
            row_top = top0 + rr * row_h
            bg = RGBColor(0x1A, 0x10, 0x40) if rr % 2 == 0 else RGBColor(0x22, 0x15, 0x55)
            _rect(slide, lft, row_top, Inches(3.5), row_h, bg)
            _rect(slide, lft + Inches(3.5), row_top, Inches(2.8), row_h, bg)
            util = float(row["utilization_pct"])
            uc   = RED if util > 100 else (YELLOW if util < 50 else GREEN)
            _txt(slide, row["month_label"], lft + Inches(0.1), row_top + Inches(0.1), Inches(3.3), row_h, size=11)
            _txt(slide, f"{util:.0f}%", lft + Inches(3.6), row_top + Inches(0.1), Inches(2.5), row_h, size=11, bold=True, color=uc)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─── Product 360 deck ─────────────────────────────────────────────────────────

def product_deck(
    prod: pd.Series,
    top_customers: pd.DataFrame,
    mom_df: pd.DataFrame,
    regional_df: pd.DataFrame,
    adoption_pct: float,
    avg_util: float,
) -> bytes:
    prs = _prs()
    today = date.today().strftime("%B %d, %Y")

    # Slide 1: Cover
    _cover(prs, str(prod["name"]), f"{prod['category']}  ·  {prod['unit']}  ·  ${float(prod['list_price_per_unit']):.3f}/unit", today)

    # Slide 2: Product snapshot KPIs
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, Inches(0), Inches(0), W, Inches(1.1), RGBColor(0x2D, 0x1F, 0x6E))
    _txt(slide, f"Product Snapshot — {prod['name']}", Inches(0.4), Inches(0.22),
         Inches(10), Inches(0.65), size=22, bold=True)

    kpis = [
        ("Category",        str(prod.get("category", "—"))),
        ("Unit",            str(prod.get("unit", "—"))),
        ("List Price",      f"${float(prod.get('list_price_per_unit', 0)):.3f}/unit"),
        ("Adoption",        f"{adoption_pct:.0f}%"),
        ("Avg Utilization", f"{avg_util:.0f}%"),
    ]
    for i, (lbl, val) in enumerate(kpis):
        _kpi_box(slide, lbl, val, Inches(0.4 + i * 2.55), Inches(1.4))

    # Slide 3: Top customers table
    if not top_customers.empty:
        slide = _blank(prs)
        _bg(slide)
        _rect(slide, Inches(0), Inches(0), W, Inches(1.1), RGBColor(0x2D, 0x1F, 0x6E))
        _txt(slide, "Top Customers", Inches(0.4), Inches(0.22), Inches(10), Inches(0.65), size=22, bold=True)

        headers = ["Customer", "Region", "Plan", "MRR", "Usage", "Utilization"]
        col_w   = [Inches(2.8), Inches(1.4), Inches(1.2), Inches(1.6), Inches(2.2), Inches(1.8)]
        row_h   = Inches(0.48)
        top0    = Inches(1.3)
        lefts   = [Inches(0.3)]
        for w in col_w[:-1]:
            lefts.append(lefts[-1] + w)

        for ci, (h, lft, ww) in enumerate(zip(headers, lefts, col_w)):
            _rect(slide, lft, top0, ww, row_h, RGBColor(0x3B, 0x2A, 0x8A))
            _txt(slide, h, lft + Inches(0.08), top0 + Inches(0.1), ww, row_h, size=10, bold=True, color=GREY)

        for ri, (_, row) in enumerate(top_customers.head(10).iterrows()):
            row_top = top0 + row_h * (ri + 1)
            bg = RGBColor(0x1A, 0x10, 0x40) if ri % 2 == 0 else RGBColor(0x22, 0x15, 0x55)
            for lft, ww in zip(lefts, col_w):
                _rect(slide, lft, row_top, ww, row_h, bg)
            util = float(row.get("utilization_pct", 0))
            util_color = RED if util > 100 else (YELLOW if util < 50 else GREEN)
            unit_str = str(row.get("unit", ""))
            vals = [
                str(row.get("name", "")),
                str(row.get("region", "")),
                str(row.get("plan_tier", "")),
                f"${float(row.get('mrr_usd', 0)):,.0f}",
                f"{float(row.get('usage', 0)):,.1f} {unit_str}",
                f"{util:.0f}%",
            ]
            for vi, (val, lft, ww) in enumerate(zip(vals, lefts, col_w)):
                c = util_color if vi == 5 else WHITE
                _txt(slide, val, lft + Inches(0.08), row_top + Inches(0.12), ww, row_h, size=10, color=c)

    # Slide 4: Regional adoption
    if not regional_df.empty:
        slide = _blank(prs)
        _bg(slide)
        _rect(slide, Inches(0), Inches(0), W, Inches(1.1), RGBColor(0x2D, 0x1F, 0x6E))
        _txt(slide, "Regional Adoption", Inches(0.4), Inches(0.22), Inches(10), Inches(0.65), size=22, bold=True)

        headers = ["Region", "Customers", "Adoption %"]
        col_w   = [Inches(4.0), Inches(2.5), Inches(2.5)]
        row_h   = Inches(0.5)
        top0    = Inches(1.3)
        lefts   = [Inches(0.3)]
        for w in col_w[:-1]:
            lefts.append(lefts[-1] + w)

        for ci, (h, lft, ww) in enumerate(zip(headers, lefts, col_w)):
            _rect(slide, lft, top0, ww, row_h, RGBColor(0x3B, 0x2A, 0x8A))
            _txt(slide, h, lft + Inches(0.08), top0 + Inches(0.1), ww, row_h, size=10, bold=True, color=GREY)

        for ri, (_, row) in enumerate(regional_df.iterrows()):
            row_top = top0 + row_h * (ri + 1)
            bg = RGBColor(0x1A, 0x10, 0x40) if ri % 2 == 0 else RGBColor(0x22, 0x15, 0x55)
            for lft, ww in zip(lefts, col_w):
                _rect(slide, lft, row_top, ww, row_h, bg)
            vals = [
                str(row.get("region", "")),
                str(int(row.get("customers", 0))),
                f"{float(row.get('adoption_pct', 0)):.0f}%",
            ]
            for val, lft, ww in zip(vals, lefts, col_w):
                _txt(slide, val, lft + Inches(0.08), row_top + Inches(0.12), ww, row_h, size=11)

    # Slide 5: MoM trend
    if not mom_df.empty:
        slide = _blank(prs)
        _bg(slide)
        _rect(slide, Inches(0), Inches(0), W, Inches(1.1), RGBColor(0x2D, 0x1F, 0x6E))
        _txt(slide, "Month-over-Month Utilization Trend", Inches(0.4), Inches(0.22), Inches(10), Inches(0.65), size=22, bold=True)
        recent = mom_df.tail(12).copy()
        recent["month_label"] = recent["month"].dt.strftime("%b %Y")
        col_w = [Inches(2.2), Inches(2.2)]
        row_h = Inches(0.42)
        top0  = Inches(1.3)
        for ri, (_, row) in enumerate(recent.iterrows()):
            col = ri % 2
            rr  = ri // 2
            lft = Inches(0.3) + col * Inches(6.5)
            row_top = top0 + rr * row_h
            bg = RGBColor(0x1A, 0x10, 0x40) if rr % 2 == 0 else RGBColor(0x22, 0x15, 0x55)
            _rect(slide, lft, row_top, Inches(3.5), row_h, bg)
            _rect(slide, lft + Inches(3.5), row_top, Inches(2.8), row_h, bg)
            util = float(row["utilization_pct"])
            uc   = RED if util > 100 else (YELLOW if util < 50 else GREEN)
            _txt(slide, row["month_label"], lft + Inches(0.1), row_top + Inches(0.1), Inches(3.3), row_h, size=11)
            _txt(slide, f"{util:.0f}%", lft + Inches(3.6), row_top + Inches(0.1), Inches(2.5), row_h, size=11, bold=True, color=uc)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
